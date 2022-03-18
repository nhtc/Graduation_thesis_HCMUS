import os, sys
from pptx import Presentation
import pandas as pd
from vncorenlp import VnCoreNLP
from docx_utils.flatten import opc_to_flat_opc
from xml.dom import minidom
import win32com.client
import uuid
import time
import pythoncom
import csv
sys.path.insert(1, os.getcwd() + "/PreprocessingComponent")

from PreprocessingComponent.pdfminer3 import Pdf_extract

if ("PreprocessingComponent" in os.getcwd()):
    vncorenlp_file = os.getcwd() + '/VnCoreNLP/VnCoreNLP-1.1.1.jar'
else:
    vncorenlp_file = os.getcwd() + '/PreprocessingComponent/VnCoreNLP/VnCoreNLP-1.1.1.jar'
vncorenlp = VnCoreNLP(vncorenlp_file, annotators="wseg,pos,ner,parse", max_heap_size='-Xmx4g', port=6000)
# annotator = VnCoreNLP(vncorenlp_file, annotators="wseg,pos,ner,parse", max_heap_size='-Xmx2g',port=6000)

"""def previous_and_next(some_iterable):
    prevs, items, nexts = tee(some_iterable, 3)
    prevs = chain([None], prevs)
    nexts = chain(islice(nexts, 1, None), [None])
    return zip(prevs, items, nexts)


# LƯU Ý:
# input của các hàm xử lý là tên file và output của các hàm xử lý ...2txt là một list các câu.
# từng phần tử của 1 list lớn (1 câu) là 1 list nhỏ chứa các từ được phân tách nhờ VNCORE.

#### ------------------------------------các function hỗ trợ cho function docx2txt
# Get paragraph string. Input is paragraph element
def para_string(para):
    string = ""
    # if (str(para)[21:34] not in str(wp_tbl)):# and (str(para)[21:34] not in str(wp_txbx)):
    wt = para.getElementsByTagName('w:t')
    for i in range(len(wt)):
        string = string + wt[i].firstChild.data

    return string


# Get table string. Input is table element
def table_string(table):
    string = ""

    wp = table.getElementsByTagName('w:p')
    column = len(table.getElementsByTagName('w:tc')) / len(table.getElementsByTagName('w:tr'))
    c = 1
    for i in range(len(wp)):
        string = string + para_string(wp[i])
        if c % column == 0 and c != len(wp):
            string = string + '. '
        else:
            string = string + '. '
        c = c + 1
    return string


# Get all elements
def get_all_elements(lst, type_of_element):
    elements_list = []
    for i in range(len(lst)):
        Elements = lst[i].getElementsByTagName(type_of_element)
        for elm in Elements:
            elements_list.append(elm)

    return elements_list


def para2text(p):
    rs = p._element.xpath('.//w:t')
    return u" ".join([r.text for r in rs])


##--------------------------------------------------------------------



def docx2txt(docx_file_name):
    # Parse xml file
    xml_file_name = 'mydocx.xml'
    print(docx_file_name)
    opc_to_flat_opc(docx_file_name, xml_file_name)
    my_docx = minidom.parse(xml_file_name)

    # Get elements
    paragraph = my_docx.getElementsByTagName('w:p')
    table = my_docx.getElementsByTagName('w:tbl')

    # Get all w:p elements in table elements. Output is two-dimensional list
    wp_tbl = get_all_elements(table, 'w:p')

    # Get text and save to "string" variable
    para_index = 0
    tbl_index = 0
    string = ""
    count = 0
    lst_para = []  # res: list paragraph of docx
    while para_index < len(paragraph):
        if paragraph[para_index] in wp_tbl:
            # string = string + table_string(table[tbl_index])
            lst_para.append(table_string(table[tbl_index]))
            para_index += len(table[tbl_index].getElementsByTagName('w:p'))
            tbl_index += 1
        else:
            # string = string + para_string(paragraph[para_index])
            lst_para.append(para_string(paragraph[para_index]))
            para_index += 1
    for i in range(0, len(lst_para)):
        if ("\xa0" in lst_para[i]):
            lst_para[i] = lst_para[i].replace("\xa0", " ")
    os.remove("mydocx.xml")
    split_sentence = []  ## list chứa danh sách câu được tách ra. mỗi phần tử là 1 câu.
    for sentences in lst_para:
        print('sentences is')
        print(sentences)
        print('end')
        split_sentence.extend(vncorenlp.pos_tag(sentences))
        #print('POS Tagging:', vncorenlp.pos_tag(sentences))
    return split_sentence  # update: trả ra pos_tag là có gán nhãn cho tưng từ về loại từ.
    # return lst_para
    # return string # gộp các đoạn lại và trả về toàn bộ văn bản ban đầu.




# hàm chuyển đổi định dạng từ fild .doc sang .docx
def doc2docx(filename, path=os.getcwd()):
    pythoncom.CoInitialize()
    baseDir = os.path.abspath(os.getcwd())  # Starting directory for directory walk
    word = win32com.client.Dispatch("Word.application")
    file_path = os.path.join(baseDir, filename)
    file_name, file_extension = os.path.splitext(file_path)

    if "~$" not in file_name:
        if file_extension.lower() == '.doc':  #
            # docx_file = '{0}{1}'.format(file_path, 'x')
            docx_file = file_name + str(uuid.uuid4().hex[:10]).format(file_path,
                                                                      'x')  # tránh trương hợp có sẵn file .docx tước đó nên thêm phần random để tránh trùng tên
            if not os.path.isfile(docx_file):  # Skip conversion where docx file already exists

                file_path = os.path.abspath(file_path)
                docx_file = os.path.abspath(docx_file)
                try:
                    wordDoc = word.Documents.Open(file_path)
                    wordDoc.SaveAs2(docx_file, FileFormat=16)
                    wordDoc.Close()
                except Exception as e:
                    print('Failed to Convert: {0}'.format(file_path))
                    print(e)
            return docx_file + ".docx"  ## trả ra tên file đã chuyển từ doc -> docx


# hàm rút trích các câu từ file ppt
def ppt2txt(filename):
    ppt = Presentation(filename)
    sentences = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                sentences += shape.text + ". "
    with VnCoreNLP(vncorenlp_file, annotators="wseg", max_heap_size='-Xmx4g', quiet=False) as vncorenlp:
        split_sentence = vncorenlp.tokenize(sentences)
    return split_sentence


# hàm rút trich các câu từ file excel ### cần cập nhật thêm vì chưa hoàn thiện
def xlsx2txt(filename):
    data = pd.read_excel(filename, sheet_name=None, index_col=0, keep_default_na=0)
    list_sheet = []
    for key, value in data.items():
        list_sheet.append(value)
    listToStr = '.'.join(map(str, list_sheet))
    list_sen = listToStr.split('\n')
    for i in range(0, len(list_sen)):
        list_sen[i] = " ".join(list_sen[i].split())
    return list_sen


def convert2listsen(vncore_postag):
    s = ""
    res = []
    index=[]
    turn=0
    for i in vncore_postag:
        count = 0
        # for previous, item, nxt in previous_and_next(i):    
        #     turn+=1
        #     if (previous == None or nxt == None):
        #         continue
        #     if (previous[0] == "[" and item[1] == "M" and nxt[0] == "]"):
        #         index.append(turn-2)
        #         index.append(turn-1)
        #         index.append(turn)
        # if(len(index)!=0):
        #     set_index=set(index)
        #     for h,k in zip(set_index,range(len(set_index))):
        #         i.pop(h-k)
        for j in i:
            if (j[1] != "CH"):
                s += j[0] + " "
            else:
                if (j[0] == '"'):
                    count += 1
                    if (count % 2 != 0):
                        s += " " + j[0]
                    else:
                        s = s.strip()
                        s += j[0] + " "
                    continue
                if (j[0] in ['(', '[', '{']):
                    s += " " + j[0]
                else:
                    s = s.strip()
                    s += j[0] + " "
        if ("_" in s):
            s = s.replace("_", " ")
        s = s.strip()
        res.append(s)
        s = ""
    return res


# hàm dùng để đém số từ trong 1 câu
def num_of_word(list_sentences):
    num_word = []
    for i in list_sentences:
        num_word.append(len(i))
    return num_word


# Hàm chia văn bản thành các đoạn
# Input:
#   + text: kiểu string, văn bản cần chia
# Output:
#   + list_para: list các đoạn

def split_text(text):
    list_para = text.split("\n\n")
    for i in range(len(list_para)):
        list_para[i] = list_para[i].replace('\n', '')
    list_para = list(filter(str.strip, list_para))
    return list_para

# Hàm đọc file pdf trả ra list các text theo các trang có áp dụng chia đoạn
# Input:
#   + file_path: Vị trị lưu file pdf cần đọc
# Output:
#   + list_para: list các văn bản đọc theo trang, 1 trang là list các đoạn văn bản của trang đó
# Ví dụ:
#   list_para[3][2]: Trang thứ 4, đoạn thứ 3

def pdf2txt_page(file_path):
    list_page = []
    with open(file_path, 'rb') as fh:
        for page in PDFPage.get_pages(fh, caching=True, check_extractable=True):
            resource_manager = PDFResourceManager()
            fake_file_handle = io.StringIO()
            converter = TextConverter(resource_manager, fake_file_handle, laparams=LAParams())
            page_interpreter = PDFPageInterpreter(resource_manager, converter)
            page_interpreter.process_page(page)
            list_page.append(fake_file_handle.getvalue())
            converter.close()
            fake_file_handle.close()
    list_para = []
    for i in range(len(list_page)):
        list_para.append(split_text(list_page[i]))
    return list_para

# Hàm đọc file pdf trả ra text và chia thành list đoạn cho text đó
# Input:
#   + file_path: Vị trị lưu file pdf cần đọc
#   + pages: nếu để trống thì đọc từ đầu đến cuối file
#           nếu có tham số là iteration thì đọc từ page a đến page b
# Output:
#   + list_para: list các đoạn của văn bản
# Ví dụ:
#   pdf2txt("sample.pdf"): Đọc toàn bộ văn bản
#   pdf2txt("sample.pdf", pages = range(1,6)): Đọc văn bản từ trang 2 đến trang 6
#   list_para[0]: Đoạn thứ 1 của văn bản
def pdf2para(file_path, pages=None):
    if not pages:
        pagenums = set()
    else:
        pagenums = set(pages)

    output = io.StringIO()
    manager = PDFResourceManager()
    converter = TextConverter(manager, output, laparams=LAParams())
    interpreter = PDFPageInterpreter(manager, converter)

    infile = open(file_path, 'rb')
    for page in PDFPage.get_pages(infile, pagenums):
        interpreter.process_page(page)
    infile.close()
    converter.close()
    text = output.getvalue()
    output.close
    list_para = split_text(text)

    return list_para # trả ra danh sách các đoạn văn bản được tách.

def pdf2txt(file_path):
    list_para = pdf2para(file_path) 
    split_sentence = []  ## list chứa danh sách câu được tách ra. mỗi phần tử là 1 câu.
    for para in list_para:
        split_sentence.extend("")  
    return split_sentence  # update: trả ra pos_tag là có gán nhãn cho tưng từ về loại từ.

# Đọc theo số trang
# Input:
#    + start_page: trang đầu tiên
#    + end_page: trang cuối cùng
#    + doc_file: file .pdf cần đọc
# Output:
#    + Tương tự output của vncorenlp.tokenize
# Ví dụ:
# Đọc từ trang 4 tới trang 7: read_pages(4, 7)
# Đọc duy nhất trang 5: read_pages(5, 5)
def read_pages(start_page, end_page, doc_file):
    words = []
    doc = pdf2txt(doc_file, range(start_page - 1, end_page))
    for para in doc:
        words.extend(vncorenlp.tokenize(para))
    return words




def preprocess(filename):
    # thực hiện đưa filename nào muốn xử lý vào biến filename bên dưới và đợi kết quả trên màn hình.
    res=[]
    #for filename in  list_filename:
    #start_time = time.time()
    #print("Thời gian bắt đầu xử lý file ",filename," la: ", start_time)
    name, file_extension = os.path.splitext(filename)
    if (file_extension.lower() == ".doc"):
        new_filename_docx = doc2docx(filename)
        a = docx2txt(new_filename_docx)
        os.remove(new_filename_docx)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".docx"):
        a = docx2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".pdf"):
        a = pdf2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)

    elif (file_extension.lower() == ".xlsx"):
        a = xlsx2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)

    elif (file_extension.lower() == ".pptx"):
        a = ppt2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)

    #res.append([os.path.basename(filename), b, num_word])  # filename, list câu. số từ của mỗi câu
    #print("Run time of file ",filename," là: --- %s seconds ---" % (time.time() - start_time))
    return os.path.basename(filename), b, num_word


# hàm này dùng dể trả ra output khác,( thêm 1 cái postag)
def preprocess_link(filename):
    # thực hiện đưa filename nào muốn xử lý vào biến filename bên dưới và đợi kết quả trên màn hình.
    res=[]
    #for filename in  list_filename:
    #start_time = time.time()
    #print("Thời gian bắt đầu xử lý file ",filename," la: ", start_time)
    name, file_extension = os.path.splitext(filename)
    if (file_extension.lower() == ".doc"):
        new_filename_docx = doc2docx(filename)
        a = docx2txt(new_filename_docx)
        os.remove(new_filename_docx)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".docx"):
        a = docx2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".pdf"):
        a = pdf2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)

    elif (file_extension.lower() == ".xlsx"):
        a = xlsx2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)

    elif (file_extension.lower() == ".pptx"):
        a = ppt2txt(filename)
        b = convert2listsen(a)  # đay là list các câu. b[0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(b)

    #res.append([os.path.basename(filename), b, num_word])  # filename, list câu. số từ của mỗi câu
    #print("Run time of file ",filename," là: --- %s seconds ---" % (time.time() - start_time))
    return a,os.path.basename(filename), b, num_word


def rtf2txt(filename):
    with open("yourfile.rtf") as infile:
        for line in infile:
            print(line)
def forfun():
    #block right after
    annotator = VnCoreNLP(vncorenlp_file, annotators="wseg,pos,ner,parse", max_heap_size='-Xmx2g',port=6000) 
    text = "Ông Nguyễn Khắc Chúc  đang làm việc tại Đại học Quốc gia Hà Nội. Bà Lan, vợ ông Chúc, cũng làm việc tại đây."
    dag ={'123':'456'}
# To perform word segmentation, POS tagging, NER and then dependency parsing
    annotated_text = annotator.annotate(text)

    hr = annotator.pos_tag(text)
    print(hr)
    # To perform word segmentation only
    word_segmented_text = annotator.tokenize(text) 
    print(annotated_text)
    print(word_segmented_text)
    #begin test

    dag ={'123':'456'}
    try:
        xml_file_name = 'mydocx.xml'
        opc_to_flat_opc('D:/project_doc.docx', xml_file_name)
        my_docx = minidom.parse(xml_file_name)
        #dag ={'123','456'}
        # Get elements
        paragraph = my_docx.getElementsByTagName('w:p')
        table = my_docx.getElementsByTagName('w:tbl')

        # Get all w:p elements in table elements. Output is two-dimensional list
        wp_tbl = get_all_elements(table, 'w:p')

        # Get text and save to "string" variable
        para_index = 0
        tbl_index = 0
        string = ""
        count = 0
        lst_para = []  # res: list paragraph of docx
        while para_index < len(paragraph):
            if paragraph[para_index] in wp_tbl:
                # string = string + table_string(table[tbl_index])
                lst_para.append(table_string(table[tbl_index]))
                para_index += len(table[tbl_index].getElementsByTagName('w:p'))
                tbl_index += 1
            else:
                # string = string + para_string(paragraph[para_index])
                lst_para.append(para_string(paragraph[para_index]))
                para_index += 1
        for i in range(0, len(lst_para)):
            if ("\xa0" in lst_para[i]):
                lst_para[i] = lst_para[i].replace("\xa0", " ")
        os.remove("mydocx.xml")
        split_sentence = []  ## list chứa danh sách câu được tách ra. mỗi phần tử là 1 câu.

        for sentences in lst_para:
            print('sentences is')
            print(sentences)
            print('end') 
            h = vncorenlp.pos_tag(sentences)
            split_sentence.extend(h)
            print('POS Tagging:', h)
        return Response(data =dag)
    except Exception as e:
        print(e)
        return Response(data =dag)
if __name__ == '__main__':
    filename = "docFile_test/bacho2.docx"
    a= preprocess(filename)
    print(a)
    annotator = VnCoreNLP(vncorenlp_file, annotators="wseg,pos,ner,parse", max_heap_size='-Xmx2g') 
    text = "Ông Nguyễn Khắc Chúc  đang làm việc tại Đại học Quốc gia Hà Nội. Bà Lan, vợ ông Chúc, cũng làm việc tại đây."

# To perform word segmentation, POS tagging, NER and then dependency parsing
    #annotated_text = annotator.annotate(text)

    # To perform word segmentation only
    #word_segmented_text = annotator.tokenize(text) 
    # print("Tên file là: ", a)
    # print("\n Danh sách các câu của file là: ", b)
    # print("\n Danh sách số từ của file là: ", c)


vncorenlp.close()"""


# hàm dùng để lặp 1 lúc 3 phần tử dùng để xóa các hyber text

# hàm dùng để lặp 1 lúc 3 phần tử dùng để xóa các hyber text
# def previous_and_next(some_iterable):
#     prevs, items, nexts = tee(some_iterable, 3)
#     prevs = chain([None], prevs)
#     nexts = chain(islice(nexts, 1, None), [None])
#     return zip(prevs, items, nexts)


# LƯU Ý:
# input của các hàm xử lý là tên file và output của các hàm xử lý ...2txt là một list các câu.
# từng phần tử của 1 list lớn (1 câu) là 1 list nhỏ chứa các từ được phân tách nhờ VNCORE.

# ------------------------------------các function hỗ trợ cho function docx2txt---------------------------
# Get paragraph string. Input is paragraph element
def para_string(para):
    string = ""
    # if (str(para)[21:34] not in str(wp_tbl)):# and (str(para)[21:34] not in str(wp_txbx)):
    wt = para.getElementsByTagName('w:t')
    for i in range(len(wt)):
        string = string + wt[i].firstChild.data

    return string


# Get table string. Input is table element
def table_string(table):
    string = ""

    wp = table.getElementsByTagName('w:p')
    column = len(table.getElementsByTagName('w:tc')) / len(table.getElementsByTagName('w:tr'))
    c = 1
    for i in range(len(wp)):
        string = string + para_string(wp[i])
        if c % column == 0 and c != len(wp):
            string = string + '. '
        else:
            string = string + '. '
        c = c + 1
    return string


# Get all elements
def get_all_elements(lst, type_of_element):
    elements_list = []
    for i in range(len(lst)):
        Elements = lst[i].getElementsByTagName(type_of_element)
        for elm in Elements:
            elements_list.append(elm)

    return elements_list


def para2text(p):
    rs = p._element.xpath('.//w:t')
    return u" ".join([r.text for r in rs])


# -------------------------------------các function hỗ trợ cho pdf2txt ------------------------------


##--------------------------------------------------------------------
# Hàm này để dùng thư viện của python xử lý file docx tự động.
# Yếu điểm là không xử lý triệt để các thứ tự các đoạn văn và phần table của văn bản


# hàm dùng để xử lý file docx sang text.
# + input là: file_path: đường dẫn đến văn bản cần xử lý
# + output là: list pos_tag của văn bản --> mục đích là biết được các câu và loại từ của từng từ trong câu.
def docx2txt(docx_file_name):
    # Parse xml file
    xml_file_name = 'mydocx.xml'
    opc_to_flat_opc(docx_file_name, xml_file_name)
    my_docx = minidom.parse(xml_file_name)

    # Get elements
    paragraph = my_docx.getElementsByTagName('w:p')
    table = my_docx.getElementsByTagName('w:tbl')

    # Get all w:p elements in table elements. Output is two-dimensional list
    wp_tbl = get_all_elements(table, 'w:p')

    # Get text and save to "string" variable
    para_index = 0
    tbl_index = 0
    string = ""
    count = 0
    lst_para = []  # res: list paragraph of docx
    while para_index < len(paragraph):
        if paragraph[para_index] in wp_tbl:
            # string = string + table_string(table[tbl_index])
            lst_para.append(table_string(table[tbl_index]))
            para_index += len(table[tbl_index].getElementsByTagName('w:p'))
            tbl_index += 1
        else:
            # string = string + para_string(paragraph[para_index])
            lst_para.append(para_string(paragraph[para_index]))
            para_index += 1
    for i in range(0, len(lst_para)):
        if ("\xa0" in lst_para[i]):
            lst_para[i] = lst_para[i].replace("\xa0", " ")
    os.remove("mydocx.xml")
    return lst_para


# hàm chuyển đổi định dạng từ fild .doc sang .docx do không tìm được cách tối ưu khi xử lỳ file .doc
def doc2docx(filename, path=os.getcwd()):
    baseDir = os.path.abspath(os.getcwd())  # Starting directory for directory walk
    pythoncom.CoInitialize()
    word = win32com.client.Dispatch("Word.application")
    file_path = os.path.join(baseDir, filename)
    file_name, file_extension = os.path.splitext(file_path)

    if "~$" not in file_name:
        if file_extension.lower() == '.doc':  #
            # docx_file = '{0}{1}'.format(file_path, 'x')
            docx_file = file_name + str(uuid.uuid4().hex[:10]).format(file_path,
                                                                      'x')  # tránh trương hợp có sẵn file .docx tước đó nên thêm phần random để tránh trùng tên
            if not os.path.isfile(docx_file):  # Skip conversion where docx file already exists

                file_path = os.path.abspath(file_path)
                docx_file = os.path.abspath(docx_file)

                try:
                    wordDoc = word.Documents.Open(file_path)
                    wordDoc.SaveAs2(docx_file, FileFormat=16)
                    wordDoc.Close()
                except Exception as e:
                    print('Failed to Convert: {0}'.format(file_path))
                    print(e)
            return docx_file + ".docx"  ## trả ra tên file đã chuyển từ doc -> docx


# hàm rút trích các câu từ file ppt
def ppt2txt(filename):
    ppt = Presentation(filename)
    res = []
    sentences = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                sentences += shape.text + ". "
        res.append(sentences)
    return res


# hàm rút trich các câu từ file excel ### cần cập nhật thêm vì chưa hoàn thiện
def xlsx2txt(filename):
    file_name, file_extension = os.path.splitext(filename)
    data = pd.read_excel(filename, index_col=0, keep_default_na=0)
    filename_csv = file_name + str(uuid.uuid4().hex[:10]) + ".csv"
    data.to_csv(filename_csv, encoding='utf-8')
    res = csv2txt(filename_csv)
    os.remove(filename_csv)
    return res


def csv2txt(filename):
    with open(filename, 'rt', encoding="utf-8")as f:
        res = []
        data = csv.reader(f)
        for row in data:
            for cell in row:
                res.append(cell + ". ")
    return res


# hàm dùng để convert pos_tag của câu thành câu hoàn chỉnh và xử lý.
def convert2listsentence(vncore_postag):
    s = ""
    res = []
    index = []
    turn = 0
    for i in vncore_postag:
        count = 0
        # for previous, item, nxt in previous_and_next(i):
        #     turn += 1
        #     if (previous == None or nxt == None):
        #         continue
        #     if (previous[0] == "[" and item[1] == "M" and nxt[0] == "]"):
        #         index.append(turn - 2)
        #         index.append(turn - 1)
        #         index.append(turn)
        # if (len(index) != 0):
        #     set_index = set(index)
        #     for h, k in zip(set_index, range(len(set_index))):
        #         i.pop(h - k)
        for j in i:
            # if (j[0] == "."): # xóa dấu chấm trong trường hợp dấu chấm lỗi bị lặp.
            #     i.remove(j)
            if (j[1] != "CH"):
                s += j[0] + " "
            else:
                if (j[0] == '"'):
                    count += 1
                    if (count % 2 != 0):
                        s += " " + j[0]
                    else:
                        s = s.strip()
                        s += j[0] + " "
                    continue
                if (j[0] in ['(', '[', '{']):
                    s += " " + j[0]
                else:
                    s = s.strip()
                    s += j[0] + " "
        if ("_" in s):
            s = s.replace("_", " ")
        s = s.strip()
        if (len(s.split(" ")) > 7):
            res.append(s)
            s = ""
    return res


# hàm dùng để đém số từ trong 1 câu
def num_of_word(list_sentences):
    num_word = []
    for i in list_sentences:
        num_word.append(len(i))
    return num_word


# hàm dùng để tách từ các đoạn văn para sang pos_tag
def list_para2txt(list_para):
    split_sentence = []  ## list chứa danh sách câu được tách ra. mỗi phần tử là 1 câu.
    for para in list_para:
        split_sentence.extend(vncorenlp.pos_tag(para))
    return split_sentence  # update: trả ra pos_tag là có gán nhãn cho tưng từ về loại từ.


def preprocess(filename):
    # thực hiện đưa filename nào muốn xử lý vào biến filename bên dưới và đợi kết quả trên màn hình.
    res = []
    # for filename in  list_filename:
    # start_time = time.time()
    # print("Thời gian bắt đầu xử lý file ",filename," la: ", start_time)'
    list_sentenc = ""
    name, file_extension = os.path.splitext(filename)
    if (file_extension.lower() not in [".doc", ".docx", ".pdf", ".xlsx", ".csv", ".pptx", ".txt"]):
        raise TypeError("Wrong type document file")
    if (file_extension.lower() == ".doc"):
        new_filename_docx = doc2docx(filename)
        list_para = docx2txt(new_filename_docx)  # list para: ds các
        pos_tag = list_para2txt(list_para)  # postag của đoạn
        os.remove(new_filename_docx)  # xóa file docx tạm
        list_sentence = convert2listsentence(
            pos_tag)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".docx"):
        list_para = docx2txt(filename)  # list para: ds các
        pos_tag = list_para2txt(list_para)  # postag của đoạn
        list_sentence = convert2listsentence(
            pos_tag)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".pdf"):
        list_para = Pdf_extract.pdf2txt(filename)  # list para: ds các
        pos_tag = list_para2txt(list_para)  # postag của đoạn
        list_sentence = convert2listsentence(
            pos_tag)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".csv"):
        list_sentence = csv2txt(
            filename)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".xlsx"):
        list_sentence = xlsx2txt(filename)  # list sentence: ds các câu
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".pptx"):
        list_para = ppt2txt(filename)  # list para: ds các
        pos_tag = list_para2txt(list_para)  # postag của đoạn
        list_sentence = convert2listsentence(
            pos_tag)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".txt"):
        f = open("docFile_test/sample.txt", "r", encoding="utf8")
        text = f.read()
        list_para = text.split("\n")  # list para: ds các
        pos_tag = list_para2txt(list_para)  # postag của đoạn
        list_sentence = convert2listsentence(
            pos_tag)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....
    # res.append([os.path.basename(filename), b, num_word])  # filename, list câu. số từ của mỗi câu
    # print("Run time of file ",filename," là: --- %s seconds ---" % (time.time() - start_time))
    return os.path.basename(filename), list_sentence, num_word


# hàm này dùng dể trả ra output khác,( thêm 1 cái postag)
def preprocess_link(filename):
    # thực hiện đưa filename nào muốn xử lý vào biến filename bên dưới và đợi kết quả trên màn hình.
    res = []
    # for filename in  list_filename:
    # start_time = time.time()
    # print("Thời gian bắt đầu xử lý file ",filename," la: ", start_time)'
    list_sentenc = ""
    name, file_extension = os.path.splitext(filename)
    if (file_extension.lower() not in [".doc", ".docx", ".pdf", ".xlsx", ".csv", ".pptx", ".txt"]):
        raise TypeError("Wrong type document file")
    if (file_extension.lower() == ".doc"):
        new_filename_docx = doc2docx(filename)
        list_para = docx2txt(new_filename_docx)  # list para: ds các
        pos_tag = list_para2txt(list_para)  # postag của đoạn
        os.remove(new_filename_docx)  # xóa file docx tạm
        list_sentence = convert2listsentence(
            pos_tag)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".docx"):
        list_para = docx2txt(filename)  # list para: ds các
        pos_tag = list_para2txt(list_para)  # postag của đoạn
        list_sentence = convert2listsentence(
            pos_tag)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".pdf"):
        list_para = Pdf_extract.pdf2txt(filename)  # list para: ds các
        pos_tag = list_para2txt(list_para)  # postag của đoạn
        list_sentence = convert2listsentence(
            pos_tag)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".csv"):
        list_sentence = csv2txt(
            filename)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".xlsx"):
        list_sentence = xlsx2txt(filename)  # list sentence: ds các câu
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".pptx"):
        list_para = ppt2txt(filename)  # list para: ds các
        pos_tag = list_para2txt(list_para)  # postag của đoạn
        list_sentence = convert2listsentence(
            pos_tag)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....

    elif (file_extension.lower() == ".txt"):
        f = open("docFile_test/sample.txt", "r", encoding="utf8")
        text = f.read()
        list_para = text.split("\n")  # list para: ds các
        pos_tag = list_para2txt(list_para)  # postag của đoạn
        list_sentence = convert2listsentence(
            pos_tag)  # đay là list các câu.  list_sentence [0] là câu đầu tiên, b[1],2,3... là các câu tiếp theo
        num_word = num_of_word(list_sentence)  # số từ của câu đầu tiên tương tự cho a[1],....
    # res.append([os.path.basename(filename), b, num_word])  # filename, list câu. số từ của mỗi câu
    # print("Run time of file ",filename," là: --- %s seconds ---" % (time.time() - start_time))
    return pos_tag, os.path.basename(filename), list_sentence, num_word


def rtf2txt(filename):
    with open("yourfile.rtf") as infile:
        for line in infile:
            print(line)


if __name__ == '__main__':
    list_filename = ["docFile_test/sample.doc", "docFile_test/sample.docx",
                     "docFile_test/sample.pdf", "docFile_test/sample.pptx", "docFile_test/lichthi.xlsx",
                     "docFile_test/lichthi.csv", "docFile_test/sample.txt"]
    for filename in list_filename:
        start_time = time.time()
        a, b, c = preprocess(filename)
        print("Tên file là: ", a)
        print("\n Danh sách các câu của file là: ", b)
        print("\n Danh sách số từ của file là: ", c)
        print("\n\nRun time of file ", filename, " là: --- %s seconds ---" % (time.time() - start_time))

vncorenlp.close()
