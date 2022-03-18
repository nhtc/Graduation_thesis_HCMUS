import tika
from tika import parser
import re
import os
import sys
import win32com.client as win32
from win32com.client import constants
from pptx import Presentation
import Extract_2
import pandas
import openpyxl
from vncorenlp import VnCoreNLP
import zipfile
import xml.etree.ElementTree
import tika
import re
from tika import parser
import logging
annotator = VnCoreNLP("VnCoreNLP//VnCoreNLP-1.1.1.jar", annotators="wseg", max_heap_size='-Xmx500m')

"""
link: https://stackoverflow.com/questions/22756344/how-do-i-extract-data-from-a-doc-docx-file-using-python
"""

WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
PARA = WORD_NAMESPACE + 'p'
TEXT = WORD_NAMESPACE + 't'
TABLE = WORD_NAMESPACE + 'tbl'
ROW = WORD_NAMESPACE + 'tr'
CELL = WORD_NAMESPACE + 'tc'


def docx2txt(filename):
    with zipfile.ZipFile('wdoc.docx') as docx:
        tree = xml.etree.ElementTree.XML(docx.read('word/document.xml'))
    s = ""
    for p in tree.iter(PARA):
        # print(''.join(node.text for node in p.iter(TEXT)))
        s = s.join(node.text for node in p.iter(TEXT))
    return s


# def doc2txt(filename):  # extract text from .doc file to string
#     tika.initVM()
#     parsed = parser.from_file(filename)
#     data = parsed["content"]
#     return data

def save_as_docx(path=os.getcwd()):  # convert .doc file to .docx file and read
    # Opening MS Word
    word = win32.gencache.EnsureDispatch('Word.Application')
    doc = word.Documents.Open(path)
    doc.Activate()

    # Rename path with .docx
    new_file_abs = os.path.abspath(path)
    new_file_abs = re.sub(r'\.\w+$', '.docx', new_file_abs)

    # Save and Close
    word.ActiveDocument.SaveAs(new_file_abs, FileFormat=constants.wdFormatXMLDocument)
    doc.Close(False)


def doc2docx(filename, path=os.getcwd()):
    # dùng khi chuyển đổi tất cả các file .doc có trong path
    # for root, dirs, filenames in os.walk(path):
    #     for f in filenames:
    #         filename, file_extension = os.path.splitext(f)
    #
    #         if file_extension.lower() == ".doc":
    #             file_conv = os.path.join(root, f)
    #             save_as_docx(file_conv)
    #             print("%s ==> %sx" % (file_conv, f))

    # dùng khi chuyển đổi 1 file
    name, file_extension = os.path.splitext(filename)
    if file_extension.lower() == ".doc":
        file_conv = os.path.join(path, filename)
        save_as_docx(file_conv)
        return filename + "x"

def ppt2text(filename):
    ppt = Presentation(filename)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                print(shape.text)


def xlsx2text(filename):
    # To open the workbook
    # workbook object is created
    wb_obj = openpyxl.load_workbook(filename)

    # Get workbook active sheet object
    # from the active attribute
    sheet_obj = wb_obj.active

    # Cell objects also have a row, column,
    # and coordinate attributes that provide
    # location information for the cell.

    # Note: The first row or
    # column integer is 1, not 0.

    # Cell object is created by using
    # sheet object's cell() method.
    cell_obj = sheet_obj.cell(row=1, column=1)

    # Print value of cell object
    # using the value attribute
    print(cell_obj.value)

def pdf2text(file):
    tika.initVM()
    parsed = parser.from_file(file)
    data=parsed["content"]
    list_sen = data.split('\s{4,}')
    for i in range (0, len(list_sen)):
        list_sen[i] = " ".join(list_sen[i].split())
    return annotator.tokenize(list_sen[0])

if __name__ == '__main__':
    annotator = VnCoreNLP("VnCoreNLP/VnCoreNLP-1.1.1.jar", annotators="wseg,pos,ner,parse", max_heap_size='-Xmx2g')

    filename = "related.pdf"

    name, file_extension = os.path.splitext(filename)

    if (file_extension.lower() == ".doc"):
        doc2docx(filename)
        print(docx2txt(filename+"x"))

    elif (file_extension.lower() == ".docx"):
        print(docx2txt(filename))
    elif(file_extension.lower() == ".pdf"):
        print(pdf2text(filename))
    elif (file_extension.lower() == ".xlxs"):
        pass
    elif (file_extension.lower() == ".pptx"):
        pass
